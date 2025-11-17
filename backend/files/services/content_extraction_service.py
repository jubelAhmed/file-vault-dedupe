"""
Content Extraction Service

Extracts text content from various file types for search indexing.
Supports:
- Text files: TXT, CSV, XML, HTML, RTF, JSON, YAML
- Documents: PDF, DOC, DOCX, ODT
- Spreadsheets: XLS, XLSX, ODS
- Presentations: PPT, PPTX, ODP
"""
import os
import logging
from typing import Optional, Set
from django.core.files.base import File as DjangoFile
from django.conf import settings
import magic

logger = logging.getLogger(__name__)


class ContentExtractionService:
    """Service for extracting searchable text content from files."""
    
    # Supported file types for content extraction
    SUPPORTED_MIME_TYPES = {
        # Text files
        'text/plain',
        'text/csv',
        'text/xml',
        'text/html',
        'text/rtf',
        'application/json',
        'application/xml',
        'application/yaml',
        'application/x-yaml',
        
        # PDF
        'application/pdf',
        
        # Word documents
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # DOCX
        'application/msword',  # DOC
        'application/vnd.oasis.opendocument.text',  # ODT
        
        # Spreadsheets
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',  # XLSX
        'application/vnd.ms-excel',  # XLS
        'application/vnd.oasis.opendocument.spreadsheet',  # ODS
        
        # Presentations
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',  # PPTX
        'application/vnd.ms-powerpoint',  # PPT
        'application/vnd.oasis.opendocument.presentation',  # ODP
    }
    
    @staticmethod
    def extract_text(file_path: str, mime_type: str) -> Optional[str]:
        """
        Extract text content from a file based on its MIME type.
        
        Args:
            file_path: Path to the file
            mime_type: MIME type of the file
            
        Returns:
            Extracted text content or None if extraction fails/unsupported
        """
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return None
            
        if mime_type not in ContentExtractionService.SUPPORTED_MIME_TYPES:
            logger.info(f"Unsupported MIME type for content extraction: {mime_type}")
            return None
        
        try:
            # Plain text files (including JSON, XML, YAML, HTML, RTF)
            if mime_type.startswith('text/') or mime_type in ['application/json', 'application/xml', 'application/yaml', 'application/x-yaml']:
                return ContentExtractionService._extract_text_file(file_path)
            
            # PDF files
            elif mime_type == 'application/pdf':
                return ContentExtractionService._extract_pdf(file_path)
            
            # Word documents
            elif 'wordprocessingml' in mime_type or mime_type == 'application/msword' or mime_type == 'application/vnd.oasis.opendocument.text':
                return ContentExtractionService._extract_docx(file_path)
            
            # Excel/Spreadsheet files
            elif 'spreadsheetml' in mime_type or mime_type == 'application/vnd.ms-excel' or mime_type == 'application/vnd.oasis.opendocument.spreadsheet':
                return ContentExtractionService._extract_xlsx(file_path)
            
            # PowerPoint/Presentation files
            elif 'presentationml' in mime_type or mime_type == 'application/vnd.ms-powerpoint' or mime_type == 'application/vnd.oasis.opendocument.presentation':
                return ContentExtractionService._extract_pptx(file_path)
            
            else:
                logger.warning(f"No extraction handler for MIME type: {mime_type}")
                return None
                
        except Exception as e:
            logger.error(f"Error extracting content from {file_path}: {str(e)}")
            return None
    
    @staticmethod
    def _extract_text_file(file_path: str) -> Optional[str]:
        """Extract content from plain text files."""
        try:
            # Try multiple encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                        logger.info(f"Successfully extracted text file with {encoding} encoding")
                        return content
                except UnicodeDecodeError:
                    continue
            
            logger.warning(f"Could not decode text file with any supported encoding: {file_path}")
            return None
            
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {str(e)}")
            return None
    
    @staticmethod
    def _extract_pdf(file_path: str) -> Optional[str]:
        """Extract text from PDF files."""
        try:
            import PyPDF2
            
            text_content = []
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text:
                        text_content.append(text)
            
            result = '\n'.join(text_content)
            logger.info(f"Extracted {len(result)} characters from PDF")
            return result if result else None
            
        except ImportError:
            logger.error("PyPDF2 not installed. Cannot extract PDF content.")
            return None
        except Exception as e:
            logger.error(f"Error extracting PDF {file_path}: {str(e)}")
            return None
    
    @staticmethod
    def _extract_docx(file_path: str) -> Optional[str]:
        """Extract text from DOCX files."""
        try:
            import docx
            
            doc = docx.Document(file_path)
            text_content = []
            
            # Extract from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text:
                    text_content.append(paragraph.text)
            
            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            text_content.append(cell.text)
            
            result = '\n'.join(text_content)
            logger.info(f"Extracted {len(result)} characters from DOCX")
            return result if result else None
            
        except ImportError:
            logger.error("python-docx not installed. Cannot extract DOCX content.")
            return None
        except Exception as e:
            logger.error(f"Error extracting DOCX {file_path}: {str(e)}")
            return None
    
    @staticmethod
    def _extract_xlsx(file_path: str) -> Optional[str]:
        """Extract text from Excel files."""
        try:
            import openpyxl
            
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' '.join([str(cell) for cell in row if cell is not None])
                    if row_text:
                        text_content.append(row_text)
            
            result = '\n'.join(text_content)
            logger.info(f"Extracted {len(result)} characters from XLSX")
            return result if result else None
            
        except ImportError:
            logger.error("openpyxl not installed. Cannot extract XLSX content.")
            return None
        except Exception as e:
            logger.error(f"Error extracting XLSX {file_path}: {str(e)}")
            return None
    
    @staticmethod
    def _extract_pptx(file_path: str) -> Optional[str]:
        """Extract text from PowerPoint presentations."""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_content = []
            
            # Extract from slides
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content.append(shape.text)
            
            result = '\n'.join(text_content)
            logger.info(f"Extracted {len(result)} characters from PPTX")
            return result if result else None
            
        except ImportError:
            logger.warning("python-pptx not installed. Cannot extract PPTX content.")
            return None
        except Exception as e:
            logger.error(f"Error extracting PPTX {file_path}: {str(e)}")
            return None
    
    @staticmethod
    def _extract_image_ocr(file_path: str) -> Optional[str]:
        """Extract text from images using OCR (Tesseract)."""
        try:
            import pytesseract
            from PIL import Image
            
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            
            logger.info(f"Extracted {len(text)} characters from image via OCR")
            return text if text.strip() else None
            
        except ImportError:
            logger.warning("pytesseract or Pillow not installed. Cannot perform OCR.")
            return None
        except Exception as e:
            logger.error(f"Error performing OCR on {file_path}: {str(e)}")
            return None
    
    @staticmethod
    def is_supported_file_type(mime_type: str) -> bool:
        """
        Check if a file type is supported for content extraction.
        
        Args:
            mime_type: MIME type to check
            
        Returns:
            True if supported, False otherwise
        """
        return mime_type in ContentExtractionService.SUPPORTED_MIME_TYPES

